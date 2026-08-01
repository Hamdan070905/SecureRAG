import easyocr
import re
import json
from functools import lru_cache
from typing import Optional
from langchain_text_splitters import RecursiveCharacterTextSplitter
import fitz
import os
import sys
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')
import time
import hashlib
import tempfile
import shutil
import chromadb
from sentence_transformers import SentenceTransformer, CrossEncoder
import fitz  # PyMuPDF
import pandas as pd
from groq import Groq
from dotenv import load_dotenv
import numpy as np
from rank_bm25 import BM25Okapi
import threading
from PIL import Image, ImageEnhance
import torch
torch.set_num_threads(4)

db_lock = threading.Lock()      # protects ChromaDB collection writes
embed_lock = threading.Lock()   # protects the shared SentenceTransformer model
ocr_lock = threading.Lock()     # protects the shared EasyOCR model

pending_image_captions = []
caption_queue_lock = threading.Lock()

def sanitize_query(query: str) -> str:
    """Basic defense against prompt injection overrides"""
    dangerous_phrases = ["ignore previous instructions", "system prompt", "you are now free", "disregard all rules"]
    q_lower = query.lower()
    for phrase in dangerous_phrases:
        if phrase in q_lower:
            return "What is the summary of the authorized documents?"
    return query

def preprocess_image(image_path: str) -> str:
    try:
        img = Image.open(image_path)
        img = img.convert('L')  # Grayscale
        # Resize to 2x using Lanczos resampling to recover blurry details
        w, h = img.size
        img = img.resize((w * 2, h * 2), Image.Resampling.LANCZOS)
        # Apply sharpness and contrast enhancements
        img = ImageEnhance.Sharpness(img).enhance(2.5)
        img = ImageEnhance.Contrast(img).enhance(2.0)
        processed_path = image_path + ".processed.png"
        img.save(processed_path)
        return processed_path
    except Exception as e:
        print("PREPROCESS OCR ERROR:", e)
        return image_path

def clean_llm_response(text: str) -> str:
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL)
    text = re.sub(r"<think>.*", "", text, flags=re.DOTALL)
    return text.strip()

load_dotenv()

print("Loading embedding model (MiniLM)...")
_EMBED_DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
embedder = SentenceTransformer("paraphrase-MiniLM-L3-v2", device=_EMBED_DEVICE)

print("Loading Cross Encoder (ms-marco-MiniLM)...")
reranker = CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2", device=_EMBED_DEVICE)

print("Cross Encoder Ready!")
# PRE-WARM (eliminates first-call delay)
try:
    embedder.encode(["warmup"], batch_size=1)
except Exception:
    pass
print("Model ready!")

def compute_embeddings(texts: list, task_type: str = "document", titles: list = None, image_paths: list = None) -> list:
    use_gemini = os.getenv("USE_GEMINI_EMBEDDING", "false").lower() == "true"
    
    if use_gemini:
        import google.generativeai as genai
        import PIL.Image
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        
        results = []
        for idx, text in enumerate(texts):
            if task_type == "query":
                content = f"task: search result | query: {text}"
            else:
                title = titles[idx] if titles and idx < len(titles) else ""
                content = f"title: {title} | text: {text}"
                
            parts = [content]
            if image_paths and idx < len(image_paths) and image_paths[idx] and os.path.exists(image_paths[idx]):
                parts.append(PIL.Image.open(image_paths[idx]))
                
            res = genai.embed_content(
                model="models/text-embedding-004",
                content=parts,
                task_type="RETRIEVAL_QUERY" if task_type == "query" else "RETRIEVAL_DOCUMENT",
                output_dimensionality=768
            )
            results.append(res['embedding'])
        return results
    else:
        return embedder.encode(
            texts,
            show_progress_bar=False,
            normalize_embeddings=True
        ).tolist()

print("Loading OCR...")
ocr_reader = easyocr.Reader(['en'], gpu=torch.cuda.is_available())
print("OCR Ready!")

groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

LLM_PROVIDER = os.getenv("LLM_PROVIDER", "ollama").lower()
LLM_MODEL = os.getenv("LLM_MODEL", "")

def set_provider(provider: str, model: str):
    global LLM_PROVIDER, LLM_MODEL
    LLM_PROVIDER = provider.lower()
    LLM_MODEL = model


def _call_llm_raw(prompt: str, max_tokens: int = 1000, temperature: float = 0.1) -> str:
    if LLM_PROVIDER == "ollama":
        from openai import OpenAI
        client = OpenAI(api_key="ollama", base_url=os.getenv("OLLAMA_BASE_URL", "http://localhost:11434/v1"))
        r = client.chat.completions.create(model=LLM_MODEL or "qwen3:8b",
            messages=[{"role": "user", "content": prompt}], max_tokens=max_tokens, temperature=temperature)
        return r.choices[0].message.content

    if LLM_PROVIDER == "groq":
        r = groq_client.chat.completions.create(model=LLM_MODEL or "llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}], max_tokens=max_tokens, temperature=temperature)
        return r.choices[0].message.content

    if LLM_PROVIDER == "openrouter":
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENROUTER_API_KEY"), base_url="https://openrouter.ai/api/v1")
        r = client.chat.completions.create(model=LLM_MODEL or "meta-llama/llama-3.3-70b-instruct",
            messages=[{"role": "user", "content": prompt}], max_tokens=max_tokens, temperature=temperature)
        return r.choices[0].message.content

    if LLM_PROVIDER == "gemini":
        import google.generativeai as genai
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel(LLM_MODEL or "gemini-2.0-flash")
        r = model.generate_content(prompt, generation_config={"max_output_tokens": max_tokens, "temperature": temperature})
        return r.text

    if LLM_PROVIDER == "openai":
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        r = client.chat.completions.create(model=LLM_MODEL or "gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}], max_tokens=max_tokens, temperature=temperature)
        return r.choices[0].message.content

    if LLM_PROVIDER == "claude":
        import anthropic
        client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
        r = client.messages.create(model=LLM_MODEL or "claude-3-5-sonnet-20241022",
            max_tokens=max_tokens, messages=[{"role": "user", "content": prompt}])
        return r.content[0].text

    if LLM_PROVIDER == "qwen":
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("QWEN_API_KEY"), base_url="https://dashscope-intl.aliyuncs.com/compatible-mode/v1")
        r = client.chat.completions.create(model=LLM_MODEL or "qwen-plus",
            messages=[{"role": "user", "content": prompt}], max_tokens=max_tokens, temperature=temperature)
        return r.choices[0].message.content

    if LLM_PROVIDER == "huggingface":
        from huggingface_hub import InferenceClient
        client = InferenceClient(token=os.getenv("HF_API_KEY"))
        return client.chat_completion(
            model=LLM_MODEL or "meta-llama/Llama-3.3-70B-Instruct",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=max_tokens, temperature=temperature,
        ).choices[0].message.content

    raise ValueError(f"Unsupported LLM_PROVIDER: {LLM_PROVIDER}")

def call_llm(prompt: str, max_tokens: int = 1000, temperature: float = 0.1) -> str:
    raw = _call_llm_raw(prompt, max_tokens, temperature)
    return clean_llm_response(raw)
chroma_client = chromadb.PersistentClient(path="vector_db")
collection = None
loaded_documents = {}
image_metadata = []
content_hashes = {}  # sha256(text) -> doc_name, for duplicate-content detection

# ==========================
# BM25 INDEX
# ==========================

bm25 = None
bm25_documents = []
bm25_metadata = []

def get_or_create_collection():
    global collection
    if collection is None:
        collection = chroma_client.get_or_create_collection(
            name="securerag",
            metadata={"hnsw:space": "cosine"}
        )
    return collection

# ============================================
# FILE EXTRACTORS
# ============================================

def extract_text(file_path: str, **kwargs):
    
    """Fast text extraction using PyMuPDF + OCR + Image Extraction"""

    global image_metadata

    ext = os.path.splitext(file_path)[1].lower()

    try:

        if ext == ".pdf":
            doc = fitz.open(file_path)
            os.makedirs("images/extracted_images", exist_ok=True)
            pages = []

            for page_number, page in enumerate(doc):
                page_text = page.get_text("text")

                # Extract tables using PyMuPDF if available
                table_text = ""
                try:
                    tables = page.find_tables()
                    for t_idx, table in enumerate(tables):
                        data = table.extract()
                        if data:
                            markdown_lines = []
                            for r_idx, row in enumerate(data):
                                row_cells = [str(cell or "").replace("\n", " ").strip() for cell in row]
                                markdown_lines.append("| " + " | ".join(row_cells) + " |")
                                if r_idx == 0:
                                    markdown_lines.append("| " + " | ".join(["---"] * len(row_cells)) + " |")
                            table_text += f"\n\n[Extracted Table {t_idx + 1}]:\n" + "\n".join(markdown_lines) + "\n\n"
                except Exception as e:
                    print(f"Table extraction error on page {page_number+1}: {e}")

                if table_text:
                    page_text += table_text

                # Skip blank pages with no text and no images
                if not page_text.strip() and not page.get_images():
                    continue

                if len(page_text.strip()) > 50:
                    # Skip expensive page render for pure-text pages in easy mode
                    page_has_images = bool(page.get_images())
                    page_image_name = f"{os.path.basename(file_path)}_p{page_number+1}.png"
                    page_image_path = os.path.join("images", "extracted_images", page_image_name)
                    if page_has_images or kwargs.get("ocr_quality", "easy") != "easy":
                        mat = fitz.Matrix(1.5, 1.5)
                        page_pix = page.get_pixmap(matrix=mat)
                        page_pix.save(page_image_path)
                        page_pix = None
                        img_path_stored = f"extracted_images/{page_image_name}"
                    else:
                        img_path_stored = None  # skip render — no PNG created

                    with db_lock:
                        image_metadata.append({
                            "doc_name": os.path.basename(file_path),
                            "page": page_number + 1,
                            "path": img_path_stored,
                            "ocr_text": page_text.lower()
                        })

                    # Searchable text page
                    pages.append({
                        "page": page_number + 1,
                        "text": page_text
                    })
                    # Process embedded diagrams/figures if present
                    images = page.get_images(full=True)
                    for img_no, img in enumerate(images):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            if pix.width < 100 or pix.height < 100:  # skip icons/decorative
                                pix = None
                                continue
                            if pix.alpha:
                                pix = fitz.Pixmap(fitz.csRGB, pix)

                            image_name = f"{os.path.basename(file_path)}_p{page_number+1}_img{img_no}.png"
                            image_path = os.path.join("images", "extracted_images", image_name)
                            pix.save(image_path)
                            pix = None

                            from ocr_router import classify_image_for_ocr
                            img_class = classify_image_for_ocr(image_path)
                            ocr_text = ""
                            img_caption = ""

                            if img_class == "ocr" and kwargs.get("ocr_quality", "off") != "off":
                                from vision_extractor import extract_image_text as _vis_embed
                                ocr_text = _vis_embed(image_path, ocr_quality=kwargs.get("ocr_quality", "easy"))
                            elif img_class == "caption" and kwargs.get("ocr_quality", "off") == "high":
                                # Queue it for background captioning
                                with caption_queue_lock:
                                    pending_image_captions.append({
                                        "image_path": image_path,
                                        "doc_name": os.path.basename(file_path),
                                        "page": page_number + 1,
                                        "file_path": file_path
                                    })
                                img_caption = ""

                            with db_lock:
                                image_metadata.append({
                                    "doc_name": os.path.basename(file_path),
                                    "page": page_number + 1,
                                    "path": f"extracted_images/{image_name}",
                                    "ocr_text": (ocr_text or img_caption or "").lower(),
                                    "img_type": img_class,
                                })

                            combined = ocr_text or img_caption
                            if combined:
                                label = "[EMBEDDED IMAGE CONTENT]" if img_class == "ocr" else "[IMAGE DESCRIPTION]"
                                pages[-1]["text"] += f"\n{label}:\n" + combined
                        except Exception as e:
                            print(f"[PDF] Embedded image error: {e}")
                else:
                    # Scanned/image-only page — always render thumbnail
                    mat = fitz.Matrix(1.5, 1.5)
                    page_pix = page.get_pixmap(matrix=mat)
                    page_image_name = f"{os.path.basename(file_path)}_p{page_number+1}.png"
                    page_image_path = os.path.join("images", "extracted_images", page_image_name)
                    page_pix.save(page_image_path)
                    page_pix = None

                    vision_page_text = ""
                    ocr_quality = kwargs.get("ocr_quality", "off")
                    if ocr_quality != "off":
                        from vision_extractor import extract_scanned_pdf_page as _scan_extract
                        vision_page_text = _scan_extract(page_image_path)

                    with db_lock:
                        image_metadata.append({
                            "doc_name": os.path.basename(file_path),
                            "page": page_number + 1,
                            "path": f"extracted_images/{page_image_name}",
                            "ocr_text": (vision_page_text or "").lower()
                        })

                    if vision_page_text.strip():
                        pages.append({
                            "page": page_number + 1,
                            "text": (page_text + "\n" + vision_page_text).strip()
                        })

            doc.close()
            return pages, None

        elif ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"):
            from vision_extractor import extract_image_text
            text = extract_image_text(file_path, ocr_quality=kwargs.get("ocr_quality", "medium"))
            print(f"[EXTRACT] Image text length: {len(text)} chars")
            
            # Save a copy to extracted_images so it can be served visually to user
            os.makedirs("images/extracted_images", exist_ok=True)
            img_name = os.path.basename(file_path)
            dest_path = os.path.join("images", "extracted_images", img_name)
            if os.path.abspath(file_path) != os.path.abspath(dest_path):
                shutil.copyfile(file_path, dest_path)
            
            with db_lock:
                image_metadata.append({
                    "doc_name": img_name,
                    "page": 1,
                    "path": f"extracted_images/{img_name}",
                    "ocr_text": (text or "").lower()
                })
            return text, None

        elif ext == ".txt":

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="ignore"
            ) as f:

                return f.read(), None

        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            pages = []
            current_heading = "Introduction"
            current_text_parts = []
            sec_idx = 1

            for p in doc.paragraphs:
                if p.style and p.style.name.startswith("Heading"):
                    if current_text_parts:
                        pages.append({
                            "page": sec_idx,
                            "text": f"Heading: {current_heading}\n" + "\n".join(current_text_parts),
                            "heading": current_heading,
                            "section": f"Section {sec_idx}",
                            "document_type": "docx",
                            "source_type": "text"
                        })
                        sec_idx += 1
                        current_text_parts = []
                    current_heading = p.text.strip()
                elif p.text.strip():
                    current_text_parts.append(p.text.strip())

            # Tables
            for tbl in doc.tables:
                table_parts = []
                for row in tbl.rows:
                    table_parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                current_text_parts.append("\n[Table]:\n" + "\n".join(table_parts))

            # Embedded images via Vision
            from vision_extractor import extract_image_text
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_data = rel.target_part.blob
                        tmp_img = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                        tmp_img.write(img_data)
                        tmp_img.close()
                        img_text = extract_image_text(tmp_img.name)
                        if img_text.strip():
                            current_text_parts.append("[Embedded Image]:\n" + img_text)
                        os.remove(tmp_img.name)
                    except Exception as img_e:
                        print(f"[DOCX] Embedded image error: {img_e}")

            if current_text_parts or current_heading:
                pages.append({
                    "page": sec_idx,
                    "text": f"Heading: {current_heading}\n" + "\n".join(current_text_parts),
                    "heading": current_heading,
                    "section": f"Section {sec_idx}",
                    "document_type": "docx",
                    "source_type": "text"
                })

            return pages, None

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            df = df.fillna("")
            cols = list(df.columns)
            col_str = ", ".join(str(c) for c in cols)
            pages = []
            
            # Summary chunk
            summary_text = f"CSV Summary\nColumns: {col_str}\nTotal Rows: {len(df)}\nTotal Columns: {len(cols)}"
            pages.append({
                "page": 1,
                "text": summary_text,
                "sheet_name": "CSV",
                "sheet_index": 1,
                "row_number": 0,
                "column_names": col_str,
                "document_type": "csv",
                "source_type": "sheet"
            })
            
            # Row chunks
            for r_idx, row in df.iterrows():
                row_parts = [f"Row: {r_idx + 1}"]
                for col in cols:
                    row_parts.append(f"{col}: {row[col]}")
                row_text = "\n".join(row_parts)
                pages.append({
                    "page": 1,
                    "text": row_text,
                    "sheet_name": "CSV",
                    "sheet_index": 1,
                    "row_number": r_idx + 1,
                    "column_names": col_str,
                    "document_type": "csv",
                    "source_type": "sheet"
                })
            return pages, None

        elif ext in (".xlsx", ".xls"):
            xls = pd.ExcelFile(file_path)
            pages = []
            for sheet_idx, sheet_name in enumerate(xls.sheet_names, 1):
                df = pd.read_excel(xls, sheet_name=sheet_name)
                df = df.fillna("")
                cols = list(df.columns)
                col_str = ", ".join(str(c) for c in cols)
                
                # Sheet Summary
                summary_text = f"Sheet Summary\nSheet: {sheet_name}\nColumns: {col_str}\nTotal Rows: {len(df)}\nTotal Columns: {len(cols)}"
                pages.append({
                    "page": sheet_idx,
                    "text": summary_text,
                    "sheet_name": sheet_name,
                    "sheet_index": sheet_idx,
                    "row_number": 0,
                    "column_names": col_str,
                    "headers": col_str,
                    "document_type": "xlsx",
                    "source_type": "excel"
                })
                
                # Rows
                for r_idx, row in df.iterrows():
                    row_parts = [f"Sheet: {sheet_name}", f"Row: {r_idx + 1}"]
                    for col in cols:
                        row_parts.append(f"{col}: {row[col]}")
                    row_text = " | ".join(row_parts)
                    pages.append({
                        "page": sheet_idx,
                        "text": row_text,
                        "sheet_name": sheet_name,
                        "sheet_index": sheet_idx,
                        "row_number": r_idx + 1,
                        "column_names": col_str,
                        "headers": col_str,
                        "document_type": "xlsx",
                        "source_type": "excel"
                    })
            return pages, None

        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                from vision_extractor import extract_image_text
                prs = Presentation(file_path)
                pages = []
                for slide_no, slide in enumerate(prs.slides, 1):
                    slide_text_parts = []
                    slide_title = ""
                    if slide.shapes.title:
                        slide_title = slide.shapes.title.text.strip()
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())
                        if shape.has_table:
                            tbl = shape.table
                            tbl_text = []
                            for r in tbl.rows:
                                tbl_text.append(" | ".join(cell.text.strip() for cell in r.cells))
                            slide_text_parts.append("\n[Table]:\n" + "\n".join(tbl_text))
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            try:
                                tmp_img = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
                                tmp_img.write(shape.image.blob)
                                tmp_img.close()
                                img_text = extract_image_text(tmp_img.name)
                                if img_text.strip():
                                    slide_text_parts.append("[Slide Image]:\n" + img_text)
                                os.remove(tmp_img.name)
                            except Exception as img_e:
                                print(f"[PPTX] Image error: {img_e}")
                    
                    # Notes
                    if slide.has_notes_slide:
                        note = slide.notes_slide.notes_text_frame.text.strip()
                        if note:
                            slide_text_parts.append(f"[Notes]: {note}")
                    
                    full_slide_text = f"Slide {slide_no}"
                    if slide_title:
                        full_slide_text += f": {slide_title}"
                    full_slide_text += "\n" + "\n".join(slide_text_parts)
                    
                    pages.append({
                        "page": slide_no,
                        "text": full_slide_text,
                        "slide_number": slide_no,
                        "slide_title": slide_title,
                        "document_type": "pptx",
                        "source_type": "slide"
                    })
                return pages, None
            except ImportError:
                return None, "python-pptx not installed"

        elif ext in (".json",):
            import json as _json
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                data = _json.load(f)
            return _json.dumps(data, indent=2), None

        elif ext in (".xml", ".md", ".markdown"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(), None

        else:

            return None, "Unsupported file"

    except Exception as e:

        return None, str(e)


# ============================================
# SMART CHUNKING
# ============================================

def smart_chunk(text: str):
    """Enterprise chunking: smaller chunks for dense/OCR text, larger for prose."""
    # Detect if text looks like OCR output (low punctuation density)
    punct_ratio = sum(1 for c in text if c in ".!?,;") / max(len(text), 1)
    is_ocr = punct_ratio < 0.015  # OCR text tends to have fewer punctuation marks

    chunk_size = 1200 if is_ocr else 1800
    overlap = 200 if is_ocr else 300

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap,
        separators=["\n\n\n", "\n\n", "\n", ". ", "! ", "? ", "; ", " ", ""],
    )
    chunks = splitter.split_text(text)
    return [c.strip() for c in chunks if len(c.strip()) > 40]




# ============================================
# DOCUMENT MANAGEMENT
# ============================================

def process_document(
    file_path,
    doc_name,
    collection_name="General",
    ocr_quality="medium"
):
    """Process document with performance benchmarking"""

    global loaded_documents

    overall_start = time.perf_counter()

    # ----------------------------
    # Extract Text
    # ----------------------------
    extract_start = time.perf_counter()

    text, error = extract_text(file_path, ocr_quality=ocr_quality)

    if error and not text:
        return 0, 0, error, None

    # PDF returns list of pages
    if isinstance(text, list):
        page_data = text
        full_text = "\n".join(page["text"] for page in page_data)

    # TXT / DOCX / CSV return string
    else:
        full_text = text
        page_data = [
            {
                "page": 1,
                "text": full_text
            }
        ]

    # <<< MOVE THIS HERE >>>
    extract_time = time.perf_counter() - extract_start

    if file_path.lower().endswith(".pdf"):
        doc = fitz.open(file_path)
        print(f"Pages      : {len(doc)}")
        doc.close()

        print(f"Characters : {len(full_text):,}")
        print(f"Words      : {len(full_text.split()):,}")

    word_count = len(full_text.split())

    # ----------------------------
    # Chunking
    # ----------------------------
    chunk_start = time.perf_counter()

    import re
    full_text = re.sub(r'\n{3,}', '\n\n', full_text)
    full_text = re.sub(r'[ \t]+', ' ', full_text)
    full_text = full_text.strip()

    # ── DEBUG: show cleaned text length before chunking ──
    print(f"\n[PIPELINE] ── {doc_name} ──")
    print(f"[PIPELINE] Raw text length   : {len(text) if isinstance(text, str) else sum(len(p['text']) for p in page_data)} chars")
    print(f"[PIPELINE] Cleaned text length: {len(full_text)} chars")
    print(f"[PIPELINE] Word count        : {len(full_text.split())} words")

    chunks = []

    for page in page_data:
        if page.get("source_type") == "sheet" or page.get("document_type") in ("xlsx", "xls", "csv"):
            page_chunks = [page["text"]]
        else:
            page_chunks = smart_chunk(page["text"])
        for c in page_chunks:
            chunks.append({
                "text": c,
                "page": page["page"],
                "slide_number": page.get("slide_number", 0),
                "slide_title": page.get("slide_title", ""),
                "sheet_index": page.get("sheet_index", 0),
                "sheet_name": page.get("sheet_name", ""),
                "heading": page.get("heading", ""),
                "section": page.get("section", ""),
                "document_type": page.get("document_type", "pdf"),
                "source_type": page.get("source_type", "text"),
                "row_number": page.get("row_number", 0),
                "column_names": page.get("column_names", ""),
            })
    # Remove duplicate chunks
    unique_chunks = []
    seen_texts = set()
    for c in chunks:
        norm = c["text"].strip().lower()
        if norm not in seen_texts:
            seen_texts.add(norm)
            unique_chunks.append(c)
    chunks = unique_chunks

    print(f"[PIPELINE] Chunk count (deduplicated): {len(chunks)}")

    if not chunks:
        print(f"[PIPELINE] ✗ 0 chunks — reasons:")
        print(f"  cleaned text length = {len(full_text)} chars")
        print(f"  smart_chunk min length filter = 40 chars")
        if len(full_text) < 40:
            print(f"  → Extracted text too short ({len(full_text)} chars). Vision/OCR likely returned empty.")
        else:
            # Force a chunk even if text is just barely there
            chunks = [{"text": full_text[:2000], "page": 1}]
            print(f"  → Fallback: creating 1 chunk from cleaned text.")

    print(f"Processing {len(chunks)} chunks for {doc_name}...")

    chunk_time = time.perf_counter() - chunk_start

    if not chunks:
        return 0, word_count, "No text chunks extracted (Vision/OCR returned empty)", None

    # ----------------------------
    # Duplicate content detection (hash of normalized full text)
    # ----------------------------
    duplicate_of = None
    text_hash = hashlib.sha256(full_text.encode("utf-8", "ignore")).hexdigest()
    with db_lock:
        existing_owner = content_hashes.get(text_hash)
        if existing_owner and existing_owner != doc_name:
            duplicate_of = existing_owner
        content_hashes[text_hash] = doc_name

    # ----------------------------
    # Generate Document ID
    # ----------------------------
    doc_id = hashlib.md5(doc_name.encode()).hexdigest()[:8]

    col = get_or_create_collection()

    try:
        existing = col.get(where={"doc_id": doc_id})
        if existing["ids"]:
            col.delete(ids=existing["ids"])
    except:
        pass

    BATCH_SIZE = 128
    all_embeddings = []
    embed_start = time.perf_counter()

    with embed_lock:
        for i in range(0, len(chunks), BATCH_SIZE):
            batch = [c["text"] for c in chunks[i:i + BATCH_SIZE]]
            titles = [c.get("doc_name", "") for c in chunks[i:i + BATCH_SIZE]]
            batch_embeddings = compute_embeddings(batch, task_type="document", titles=titles)
            all_embeddings.extend(batch_embeddings)

        print(f"Embedded {min(len(chunks), (len(chunks) // BATCH_SIZE + 1) * BATCH_SIZE)}/{len(chunks)}")

        embed_time = time.perf_counter() - embed_start

        # ----------------------------
        # Store in ChromaDB
        # ----------------------------
        store_start = time.perf_counter()

        STORE_BATCH = 100
        stored = 0

        for i in range(0, len(chunks), STORE_BATCH):
            batch_chunks = chunks[i:i + STORE_BATCH]
            batch_embeddings = all_embeddings[i:i + STORE_BATCH]

            batch_ids = [
                f"{doc_id}_chunk_{i + j}"
                for j in range(len(batch_chunks))
            ]

            batch_meta = []

            upload_time = time.strftime("%Y-%m-%d %H:%M:%S")
            for j, chunk in enumerate(batch_chunks):
                batch_meta.append({
                    "doc_name": doc_name,
                    "doc_id": doc_id,
                    "chunk_index": i + j,
                    "page": chunk["page"],
                    "total_chunks": len(chunks),
                    "previous_chunk": max(0, i + j - 1),
                    "next_chunk": min(len(chunks) - 1, i + j + 1),
                    "upload_time": upload_time,
                    "file_hash": text_hash,
                    "file_path": file_path,  # for lazy PDF page preview
                    "chunk_id": f"{doc_id}_chunk_{i + j}",
                    "parent_chunk_id": chunk.get("parent_chunk_id", ""),
                    "source_type": chunk.get("source_type", "text"),
                    "image_path": chunk.get("image_path", ""),
                    "image_caption": chunk.get("image_caption", ""),
                    "ocr_engine_used": chunk.get("ocr_engine_used", ""),
                    "slide_number": int(chunk.get("slide_number", 0)),
                    "slide_title": str(chunk.get("slide_title", "")),
                    "sheet_index": int(chunk.get("sheet_index", 0)),
                    "sheet_name": str(chunk.get("sheet_name", "")),
                    "heading": str(chunk.get("heading", "")),
                    "section": str(chunk.get("section", "")),
                    "document_type": str(chunk.get("document_type", "pdf")),
                    "row_number": int(chunk.get("row_number", 0)),
                    "column_names": str(chunk.get("column_names", "")),
                })

            col.add(
                embeddings=batch_embeddings,
                documents=[c["text"] for c in batch_chunks],
                ids=batch_ids,
                metadatas=batch_meta,
            )
            stored += len(batch_chunks)

    # Thread-safe BM25 rebuild & Metadata Save
    with db_lock:
        global bm25, bm25_documents, bm25_metadata
        for idx, chunk in enumerate(chunks):
            bm25_documents.append(chunk["text"].split())
            bm25_metadata.append({
                "text": chunk["text"], "page": chunk["page"],
                "doc_name": doc_name, "chunk_index": idx,
                "file_path": file_path,
                "chunk_id": f"{doc_id}_chunk_{idx}",
                "parent_chunk_id": chunk.get("parent_chunk_id", ""),
                "source_type": chunk.get("source_type", "text"),
                "image_path": chunk.get("image_path", ""),
                "image_caption": chunk.get("image_caption", ""),
                "ocr_engine_used": chunk.get("ocr_engine_used", ""),
                "slide_number": int(chunk.get("slide_number", 0)),
                "slide_title": str(chunk.get("slide_title", "")),
                "sheet_index": int(chunk.get("sheet_index", 0)),
                "sheet_name": str(chunk.get("sheet_name", "")),
                "heading": str(chunk.get("heading", "")),
                "section": str(chunk.get("section", "")),
                "document_type": str(chunk.get("document_type", "pdf")),
                "row_number": int(chunk.get("row_number", 0)),
                "column_names": str(chunk.get("column_names", "")),
            })
        bm25 = BM25Okapi(bm25_documents)
        
        # Keep this INSIDE the lock, remove the one outside
        loaded_documents[doc_name] = {
            "doc_id": doc_id, "chunks": stored,
            "words": word_count, "chars": len(full_text), "file": doc_name,
            "collection": collection_name or "General",
            "duplicate_of": duplicate_of,
            "summary": None, "suggested_questions": [],
        }

    clear_retrieval_cache()

    # ----------------------------
    # Performance Calculations
    # ----------------------------
    store_time = time.perf_counter() - store_start
    total_time = time.perf_counter() - overall_start

    print("\n========== PERFORMANCE ==========")
    print(f"Extraction : {extract_time:.2f} sec")
    print(f"Chunking   : {chunk_time:.2f} sec")
    print(f"Embedding  : {embed_time:.2f} sec")
    print(f"Storage    : {store_time:.2f} sec")
    print(f"Total      : {total_time:.2f} sec")
    print("=================================\n")

    print(f"Done: {stored} chunks stored.")

    return stored, word_count, None, duplicate_of

def get_loaded_documents() -> dict:
    return loaded_documents


def lazy_render_page_preview(file_path: str, page_number: int) -> str | None:
    """
    Render a single PDF page to PNG only when first requested.
    Returns the relative path (extracted_images/xxx.png) or None.
    """
    try:
        fname = os.path.basename(file_path)
        page_image_name = f"{fname}_p{page_number}.png"
        dest = os.path.join("images", "extracted_images", page_image_name)
        if os.path.exists(dest):
            return f"extracted_images/{page_image_name}"  # cached
        if not os.path.exists(file_path):
            return None
        doc = fitz.open(file_path)
        if page_number < 1 or page_number > len(doc):
            doc.close()
            return None
        page = doc[page_number - 1]
        mat = fitz.Matrix(2.0, 2.0)  # higher DPI for evidence viewer
        pix = page.get_pixmap(matrix=mat)
        os.makedirs("images/extracted_images", exist_ok=True)
        pix.save(dest)
        pix = None
        doc.close()
        return f"extracted_images/{page_image_name}"
    except Exception as e:
        print(f"[PREVIEW] Failed to render page {page_number} of {file_path}: {e}")
        return None

def clear_all_documents():
    global collection
    global loaded_documents
    global bm25
    global bm25_documents
    global bm25_metadata
    global content_hashes

    clear_retrieval_cache()

    try:
        chroma_client.delete_collection("securerag")
    except Exception:
        pass

    collection = None
    loaded_documents = {}
    bm25 = None
    bm25_documents = []
    bm25_metadata = []
    image_metadata.clear()
    content_hashes.clear()

# ============================================
# RETRIEVAL
# ============================================

_retrieval_cache = {}
_cache_lock = threading.Lock()

def clear_retrieval_cache():
    with _cache_lock:
        _retrieval_cache.clear()
    # Also bust the lru_cache on embeddings (safe to do; it just re-encodes next call)
    _cached_query_embedding_tuple.cache_clear()

@lru_cache(maxsize=256)
def _cached_query_embedding_tuple(query: str):
    return tuple(compute_embeddings([query], task_type="query")[0])

def _cached_query_embedding(query: str):
    return list(_cached_query_embedding_tuple(query))

def retrieve_chunks(
    query,
    top_k=5,
    doc_filter=None
):
    global _retrieval_cache

    # ─ 1. Resolve live docs ──────────────────────────────────
    with db_lock:
        live_docs = set(loaded_documents.keys())

    if not live_docs:
        print("[RETRIEVAL] No documents loaded — returning []")
        return []

    # ─ 2. Resolve filter_set ─────────────────────────────
    if doc_filter and doc_filter != "All Documents":
        requested = {doc_filter} if isinstance(doc_filter, str) else set(doc_filter)
        filter_set = requested & live_docs
    else:
        filter_set = live_docs

    if not filter_set:
        print(f"[RETRIEVAL] Requested doc not in loaded_documents: {doc_filter}")
        return []

    filter_list = sorted(filter_set)
    print(f"[RETRIEVAL] Query='{query[:60]}' | Filter={filter_list} | top_k={top_k}")

    cache_key = (query, top_k, tuple(filter_list))
    with _cache_lock:
        if cache_key in _retrieval_cache:
            print("[RETRIEVAL] Cache hit")
            return _retrieval_cache[cache_key]

    col = get_or_create_collection()
    if col.count() == 0:
        return []

    query_embedding = _cached_query_embedding(query)

    # ─ 3. Chroma where-filter ───────────────────────────
    where_filter = (
        {"doc_name": {"$in": filter_list}}
        if len(filter_list) > 1
        else {"doc_name": filter_list[0]}
    )
    print(f"[RETRIEVAL] Chroma where={where_filter}")

    try:
        # ─ VECTOR SEARCH ──────────────────────────────────
        results = col.query(
            query_embeddings=[query_embedding],
            n_results=min(top_k * 2, col.count()),  # over-fetch for reranking
            where=where_filter,
            include=["documents", "metadatas", "distances"]
        )

        chunks = []
        seen = set()

        # Parse query for explicit page, slide, sheet, section, or figure targets
        meta_target = {}
        query_lower = query.lower()
        
        slide_match = re.search(r'\bslide\s*(\d+)\b', query_lower)
        if slide_match:
            meta_target["slide_number"] = int(slide_match.group(1))
            
        page_match = re.search(r'\bpage\s*(\d+)\b', query_lower)
        if page_match:
            meta_target["page"] = int(page_match.group(1))
            
        sheet_match = re.search(r'\bsheet\s*(\d+)\b', query_lower)
        if sheet_match:
            meta_target["sheet_index"] = int(sheet_match.group(1))
            
        sheet_name_match = re.search(r'\b(?:worksheet|sheet)\s*([a-zA-Z0-9_\-]+)\b', query_lower)
        if sheet_name_match and not sheet_match:
            meta_target["sheet_name"] = sheet_name_match.group(1)
            
        section_match = re.search(r'\b(?:section|chapter)\s*(\d+)\b', query_lower)
        if section_match:
            meta_target["section"] = int(section_match.group(1))

        row_match = re.search(r'\brow\s*(\d+)\b', query_lower)
        if row_match:
            meta_target["row_number"] = int(row_match.group(1))

        for doc, meta, dist in zip(
            results["documents"][0],
            results["metadatas"][0],
            results["distances"][0]
        ):
            dn = meta.get("doc_name", "Unknown")
            if dn not in filter_set:
                print(f"[RETRIEVAL] REJECTED stale chunk from '{dn}'")
                continue
            chunks.append({
                "text": doc, "doc_name": dn,
                "chunk_index": meta.get("chunk_index", 0),
                "page": meta.get("page", 1),
                "file_path": meta.get("file_path"),
                "source_type": meta.get("source_type", "text"),
                "image_path": meta.get("image_path", ""),
                "image_caption": meta.get("image_caption", ""),
                "chunk_id": meta.get("chunk_id", ""),
                "parent_chunk_id": meta.get("parent_chunk_id", ""),
                "relevance": round((1 - dist) * 100, 1),
                # Structural metadata
                "slide_number": meta.get("slide_number", 0),
                "slide_title": meta.get("slide_title", ""),
                "sheet_index": meta.get("sheet_index", 0),
                "sheet_name": meta.get("sheet_name", ""),
                "heading": meta.get("heading", ""),
                "section": meta.get("section", ""),
                "document_type": meta.get("document_type", "pdf"),
                "row_number": meta.get("row_number", 0),
                "column_names": meta.get("column_names", ""),
            })
            seen.add(doc)

        # Explicitly pull in targeted metadata chunks to guarantee presence in retrieval pool
        if meta_target:
            with db_lock:
                bm25_meta_snapshot = list(bm25_metadata) if bm25_metadata else []
            for chunk in bm25_meta_snapshot:
                if chunk["doc_name"] in filter_set and chunk["text"] not in seen:
                    match = False
                    if "slide_number" in meta_target and chunk.get("slide_number") == meta_target["slide_number"]:
                        match = True
                    if "page" in meta_target and chunk.get("page") == meta_target["page"]:
                        match = True
                    if "sheet_index" in meta_target and chunk.get("sheet_index") == meta_target["sheet_index"]:
                        match = True
                    if "sheet_name" in meta_target and str(chunk.get("sheet_name")).lower() == str(meta_target["sheet_name"]).lower():
                        match = True
                    if "section" in meta_target and str(meta_target["section"]) in str(chunk.get("section")).lower():
                        match = True
                    if "row_number" in meta_target and chunk.get("row_number") == meta_target["row_number"]:
                        match = True
                        
                    if match:
                        chunks.append({
                            "text": chunk["text"], "doc_name": chunk["doc_name"],
                            "chunk_index": chunk["chunk_index"], "page": chunk["page"],
                            "file_path": chunk.get("file_path"),
                            "source_type": chunk.get("source_type", "text"),
                            "image_path": chunk.get("image_path", ""),
                            "image_caption": chunk.get("image_caption", ""),
                            "chunk_id": chunk.get("chunk_id", ""),
                            "parent_chunk_id": chunk.get("parent_chunk_id", ""),
                            "relevance": 95.0,
                            "slide_number": chunk.get("slide_number", 0),
                            "slide_title": chunk.get("slide_title", ""),
                            "sheet_index": chunk.get("sheet_index", 0),
                            "sheet_name": chunk.get("sheet_name", ""),
                            "heading": chunk.get("heading", ""),
                            "section": chunk.get("section", ""),
                            "document_type": chunk.get("document_type", "pdf"),
                            "row_number": chunk.get("row_number", 0),
                            "column_names": chunk.get("column_names", ""),
                        })
                        seen.add(chunk["text"])

        # ─ BM25 SEARCH ───────────────────────────────────
        with db_lock:
            bm25_docs_snapshot = list(bm25_documents) if bm25_documents else []
            bm25_meta_snapshot = list(bm25_metadata) if bm25_metadata else []
            bm25_snapshot = bm25

        if bm25_snapshot is not None:
            tokenized_query = query.lower().split()
            scores = bm25_snapshot.get_scores(tokenized_query)
            filtered_indices = [
                idx for idx, m in enumerate(bm25_meta_snapshot)
                if m["doc_name"] in filter_set
            ]
            filtered_indices.sort(key=lambda i: scores[i], reverse=True)
            for idx in filtered_indices[:top_k * 2]:
                chunk = bm25_meta_snapshot[idx]
                if chunk["text"] in seen:
                    continue
                chunks.append({
                    "text": chunk["text"], "doc_name": chunk["doc_name"],
                    "chunk_index": chunk["chunk_index"], "page": chunk["page"],
                    "file_path": chunk.get("file_path"),
                    "source_type": chunk.get("source_type", "text"),
                    "image_path": chunk.get("image_path", ""),
                    "image_caption": chunk.get("image_caption", ""),
                    "chunk_id": chunk.get("chunk_id", ""),
                    "parent_chunk_id": chunk.get("parent_chunk_id", ""),
                    "relevance": round(float(scores[idx]), 1),
                    # Structural metadata
                    "slide_number": chunk.get("slide_number", 0),
                    "slide_title": chunk.get("slide_title", ""),
                    "sheet_index": chunk.get("sheet_index", 0),
                    "sheet_name": chunk.get("sheet_name", ""),
                    "heading": chunk.get("heading", ""),
                    "section": chunk.get("section", ""),
                    "document_type": chunk.get("document_type", "pdf"),
                    "row_number": chunk.get("row_number", 0),
                    "column_names": chunk.get("column_names", ""),
                })
                seen.add(chunk["text"])

        # ─ EXACT KEYWORD & METADATA BOOST ────────────────────────────
        for c in chunks:
            boost = 0.0
            # Target Boosting
            if "slide_number" in meta_target and c.get("slide_number") == meta_target["slide_number"]:
                boost += 45.0
            if "page" in meta_target and c.get("page") == meta_target["page"]:
                boost += 45.0
            if "sheet_index" in meta_target and c.get("sheet_index") == meta_target["sheet_index"]:
                boost += 45.0
            if "sheet_name" in meta_target and str(c.get("sheet_name")).lower() == str(meta_target["sheet_name"]).lower():
                boost += 45.0
            if "section" in meta_target and (str(meta_target["section"]) in str(c.get("section")).lower() or str(meta_target["section"]) in str(c.get("page")).lower()):
                boost += 45.0
            if "row_number" in meta_target and c.get("row_number") == meta_target["row_number"]:
                boost += 45.0
            
            # Column name presence boost
            if c.get("column_names"):
                cols_list = [col.strip().lower() for col in c["column_names"].split(",")]
                for col in cols_list:
                    if col and col in query_lower:
                        boost += 25.0
                        break
                        
            # Sheet name presence boost
            if c.get("sheet_name") and c["sheet_name"].lower() in query_lower:
                boost += 35.0
            
            # Exact keyword boost
            if query_lower in c["text"].lower():
                boost += 20.0
                
            if boost > 0:
                c["relevance"] = min(100.0, c["relevance"] + boost)

        print(f"[RETRIEVAL] Pre-rerank pool: {len(chunks)} chunks from docs: "
              f"{list({c['doc_name'] for c in chunks})}")

        # ─ RERANK + TRIM + EXPAND ────────────────────────
        chunks = rerank_chunks(query, chunks)
        chunks = chunks[:top_k]
        chunks = expand_parent_context(chunks)

        print(f"[RETRIEVAL] Final: {len(chunks)} chunks | docs: "
              f"{list({c['doc_name'] for c in chunks})}")

        with _cache_lock:
            if len(_retrieval_cache) > 200:
                _retrieval_cache.clear()
            _retrieval_cache[cache_key] = chunks

        return chunks

    except Exception as e:
        print(f"[RETRIEVAL] ERROR: {e}")
        return []
    
def expand_parent_context(chunks):

    expanded = []

    for chunk in chunks:

        index = chunk["chunk_index"]

        doc = chunk["doc_name"]

        parent = []

        for offset in [-1, 0, 1]:

            idx = index + offset

            if idx < 0:
                continue

            results = collection.get(
                where={
                    "$and": [
                         {"doc_name": doc},
                         {"chunk_index": idx}
                    ]
                }
            )

            if results["documents"]:

                parent.append(
                    results["documents"][0]
                )

        chunk["text"] = "\n\n".join(parent)

        expanded.append(chunk)

    return expanded

# ============================================
# QUERY EXPANSION
# ============================================

def expand_query(query):

    expansions = {
        "leave": [
            "annual leave",
            "vacation",
            "employee leave",
            "paid leave"
        ],

        "salary": [
            "compensation",
            "pay",
            "income",
            "wages"
        ],

        "policy": [
            "guidelines",
            "rules",
            "procedure",
            "regulation"
        ],

        "security": [
            "authentication",
            "authorization",
            "guardrails",
            "protection"
        ],

        "document": [
            "file",
            "pdf",
            "report",
            "manual"
        ]
    }

    queries = [query]

    lower = query.lower()

    for word, synonyms in expansions.items():

        if word in lower:

            queries.extend(synonyms)

    return list(set(queries))


# ============================================
# CROSS ENCODER RERANKER
# ============================================

def rerank_chunks(query, chunks):

    if len(chunks) <= 1:
        return chunks

    pairs = [
        (query, chunk["text"])
        for chunk in chunks
    ]

    scores = reranker.predict(pairs)

    for chunk, score in zip(chunks, scores):
        chunk["rerank_score"] = float(score)

    chunks.sort(
        key=lambda x: x["rerank_score"],
        reverse=True
    )

    return chunks

import re

def build_response_instruction(query: str) -> str:
    q = query.lower()

    # Dynamic word count detection
    m = re.search(
        r"(?:about|around|approximately|approx|maximum|max|less than|under|within|exactly|in)?\s*(\d+)\s*words?",
        q
    )

    if m:
        n = int(m.group(1))

        if "less than" in q or "under" in q:
            return f"Keep the answer under {n} words."

        if "maximum" in q or "max" in q:
            return f"Do not exceed {n} words."

        if "around" in q or "about" in q or "approximately" in q:
            return f"Write approximately {n} words."

        return f"Write approximately {n} words."

    if "one sentence" in q:
        return "Answer in one sentence."

    if "two sentence" in q:
        return "Answer in two sentences."

    if "one paragraph" in q:
        return "Answer in one paragraph."

    if "bullet" in q:
        return "Answer using bullet points."

    if "step-by-step" in q:
        return "Answer as step-by-step instructions."

    if any(x in q for x in ["short answer", "brief", "concise"]):
        return "Provide a concise answer."

    if any(x in q for x in ["detailed", "elaborate"]):
        return "Provide a detailed explanation."

    return "Provide a clear and well-structured response."

import json # Make sure json is imported at the top of the file if it isn't already

def evaluate_answer(query: str, answer: str, chunks: list) -> dict:
    context = "\n\n".join(c["text"][:400] for c in chunks[:5])
    prompt = f"""Rate this answer 0-100 on two axes. Respond ONLY as JSON.
Context: {context[:3000]}
Question: {query}
Answer: {answer[:1500]}
{{"faithfulness": <int, is the answer supported by context>, "groundedness": <int, does it avoid unsupported claims>}}"""
    try:
        raw = call_llm(prompt, max_tokens=60, temperature=0).strip().replace("```json", "").replace("```", "")
        return json.loads(raw)
    except Exception:
        return {"faithfulness": None, "groundedness": None}

def build_relations(doc_filter=None):
    """Fallback for knowledge graph relations."""
    return []

def query_graph(entity: str):
    """Fallback for knowledge graph query."""
    return []
# ============================================
# AI ANSWER GENERATION
# ============================================

def generate_answer(query: str, chunks: list, mode: str = "detailed") -> dict:
    """Generate grounded answer with confidence score and adaptable length based on intent"""
    
    if not chunks:
        return {
            "answer": "No documents loaded. Please upload a document first.",
            "confidence": 0,
            "sources": []
        }
    
    context_parts = []
    sources = []
    MAX_CONTEXT = 8000  # Expanded token window to capture deeper text blocks
    current_length = 0

    for chunk in chunks:
        source_type = chunk.get("source_type", "text")
        prefix = f"[Source: {chunk['doc_name']} | Page: {chunk.get('page', 1)} | Relevance: {chunk['relevance']}%"
        if source_type == "image":
            prefix += " | Type: Image Caption"
        prefix += "]\n"
        chunk_text = prefix + chunk['text']

        if current_length + len(chunk_text) > MAX_CONTEXT:
            break

        context_parts.append(chunk_text)
        current_length += len(chunk_text)

        if chunk['doc_name'] not in sources:
            sources.append(chunk['doc_name'])

    context = "\n\n---\n\n".join(context_parts)
    avg_relevance = sum(c['relevance'] for c in chunks) / len(chunks)

    # Dynamic length and format guidance based on prompt instructions
    query_lower = query.lower()
    if any(w in query_lower for w in ["short", "brief", "summary", "summarize", "in brief"]):
        length_instruction = "Provide a concise summary response in 2-4 sentences."
    elif any(w in query_lower for w in ["detail", "explain", "comprehensive", "extract the whole", "fully"]):
        length_instruction = "Provide a thorough, comprehensive, and detailed extraction/explanation drawing directly from all matching sections of the text."
    elif any(w in query_lower for w in ["bullet", "points", "list"]):
        length_instruction = "Break down the answer using clear structured bullet points."
    else:
        length_instruction = "Provide a structured, accurate, and contextually complete response based on the provided data."



    prompt = f"""You are a secure Enterprise RAG assistant. You ONLY answer from the provided document context.

STRICT RULES:
1. Answer ONLY using the Context below. Do NOT use any external knowledge.
2. {length_instruction}
3. If the answer is not in the context, reply EXACTLY: "The document does not contain this information."
4. ALWAYS cite: document name + page number for every claim.
5. Never hallucinate, speculate, or infer beyond what is written.
6. If text appears garbled/low-quality, extract what you can and note: "[Note: Source image was low quality — extracted text may be partial.]"
7. If the query is in a different language than the document, still answer in the query's language.
8. If the information comes from a chunk labeled as "Type: Image Caption", explicitly mention in your answer that it is based on an image (e.g., "According to the architectural render...").

Document Context:
{context}

User Question: {query}

Answer (with citations):"""
    try:
        answer = call_llm(
            prompt,
            max_tokens=1500,
            temperature=0.1,
        )
        
        confidence = min(98, int(avg_relevance))
        
        return {
            "answer": answer,
            "confidence": confidence,
            "sources": sources,
            "chunks_used": len(chunks),
            "avg_relevance": avg_relevance
        }
    
    except Exception as e:
        return {
            "answer": f"Error generating answer: {str(e)}",
            "confidence": 0,
            "sources": []
        }

def generate_answer_stream(query, chunks, mode="detailed"):
    if not chunks:
        def _empty():
            yield "No documents are currently loaded. Please upload a document first."
        return _empty()
    context_parts = []
    for c in chunks[:10]:
        stype = c.get("source_type", "text")
        t_label = " | Type: Image Caption" if stype == "image" else ""
        context_parts.append(f"[{c['doc_name']} | Page {c.get('page',1)}{t_label}]\n{c['text']}")
    context = "\n\n".join(context_parts)
    
    prompt = f"""You are a secure RAG assistant. Answer ONLY from the context below.
If the answer is not present, say: "The document does not contain this information."
Always cite document name and page number.
Never hallucinate.
If the information comes from a chunk labeled "Type: Image Caption", explicitly mention that the answer is based on an image.

Context:
{context}

Question: {query}
Answer:"""
    if LLM_PROVIDER == "groq":
        stream = groq_client.chat.completions.create(
            model=LLM_MODEL or "llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            stream=True,
        )
        for chunk in stream:
            delta = chunk.choices[0].delta.content
            if delta:
                yield delta
    else:
        yield call_llm(prompt, max_tokens=1500)  # non-streaming fallback

    


# ============================================
# SINGLE-DOCUMENT DELETE (for per-file remove in UI)
# ============================================

def delete_document(doc_name: str) -> bool:
    global bm25, bm25_documents, bm25_metadata, loaded_documents, content_hashes

    if doc_name not in loaded_documents:
        return False

    doc_id = loaded_documents[doc_name].get("doc_id")
    col = get_or_create_collection()

    with db_lock:
        try:
            if doc_id:
                existing = col.get(where={"doc_id": doc_id})
                if existing["ids"]:
                    col.delete(ids=existing["ids"])
        except Exception:
            pass

        keep_idx = [i for i, m in enumerate(bm25_metadata) if m["doc_name"] != doc_name]
        bm25_documents = [bm25_documents[i] for i in keep_idx]
        bm25_metadata = [bm25_metadata[i] for i in keep_idx]
        bm25 = BM25Okapi(bm25_documents) if bm25_documents else None

        image_metadata[:] = [m for m in image_metadata if m["doc_name"] != doc_name]

        # Remove content hash so the same file can be re-uploaded cleanly
        content_hashes = {h: n for h, n in content_hashes.items() if n != doc_name}

        del loaded_documents[doc_name]

    clear_retrieval_cache()

    return True


# ============================================
# AI DOCUMENT SUMMARY + SUGGESTED QUESTIONS
# ============================================

def generate_summary_and_questions(doc_name: str) -> dict:
    """Runs once after upload: short summary + 3 suggested starter questions."""
    sample_chunks = [m for m in bm25_metadata if m["doc_name"] == doc_name][:8]
    if not sample_chunks:
        return {"summary": "", "suggested_questions": []}

    context = "\n\n".join(c["text"][:600] for c in sample_chunks)
    prompt = f"""Summarize the document below in 2-3 sentences, then suggest 3 short
questions a user could ask about it. Respond ONLY as JSON:
{{"summary": "...", "suggested_questions": ["...", "...", "..."]}}

Document excerpt:
{context[:6000]}
"""
    try:
        raw = call_llm(
            prompt,
            max_tokens=400,
            temperature=0.2,
        ).strip()

        raw = raw.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)
        loaded_documents[doc_name]["summary"] = data.get("summary", "")
        loaded_documents[doc_name]["suggested_questions"] = data.get("suggested_questions", [])
        return data
    except Exception as e:
        return {"summary": "", "suggested_questions": [], "error": str(e)}


# ============================================
# AI ACTIONS: Summarize / Explain / Translate / Compare / Extract
# ============================================

ACTION_PROMPTS = {
    "summarize": "Summarize the following document content clearly and concisely.",
    "explain": "Explain the following document content in simple, plain terms as if teaching a beginner.",
    "translate": "Translate the following document content into {language}. Preserve meaning and structure.",
    "compare": "Compare and contrast the content from the different documents below. Highlight key similarities and differences.",
    "extract": "Extract all key facts, figures, dates, and action items from the following content as a structured bullet list.",
}

def generate_action_answer(action: str, chunks: list, language: str = "Spanish") -> dict:
    if action not in ACTION_PROMPTS:
        return {"answer": "Unknown action.", "confidence": 0, "sources": []}
    if not chunks:
        return {"answer": "No documents loaded.", "confidence": 0, "sources": []}

    # Extract mode: return raw text from chunks — no LLM, no paraphrasing
    if action == "extract":
        parts = [
            f"[{c['doc_name']} | Page {c.get('page', 1)}]\n{c['text']}"
            for c in chunks[:15]
        ]
        sources = list(dict.fromkeys(c["doc_name"] for c in chunks[:15]))
        return {"answer": "\n\n---\n\n".join(parts), "confidence": 95, "sources": sources}

    context = "\n\n---\n\n".join(
        f"[Source: {c['doc_name']} | Page: {c.get('page', 1)}]\n{c['text']}" for c in chunks[:10]
    )
    instruction = ACTION_PROMPTS[action].format(language=language)
    prompt = f"{instruction}\n\nContent:\n{context[:8000]}\n\nResponse:"

    try:
        answer_text = call_llm(
            prompt,
            max_tokens=1200,
            temperature=0.2,
        )
        sources = list({c["doc_name"] for c in chunks[:10]})
        return {"answer": answer_text, "confidence": 90, "sources": sources}
    except Exception as e:
        return {"answer": f"Error: {e}", "confidence": 0, "sources": []}


# ============================================
# BASIC INSIGHTS: entities, timeline, knowledge graph
# (Lightweight regex heuristics — no extra ML dependency)
# ============================================

_DATE_RE = re.compile(
    r"\b(?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4}|"
    r"(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}|"
    r"\d{4})\b"
)
_ENTITY_RE = re.compile(r"\b(?:[A-Z][a-z]+(?:\s+[A-Z][a-z]+){0,2})\b")
_STOPWORDS = {"The", "This", "That", "These", "Those", "It", "In", "On", "For", "As", "A", "An"}

def generate_insights(doc_filter: Optional[list] = None) -> dict:
    rows = bm25_metadata if not doc_filter else [m for m in bm25_metadata if m["doc_name"] in doc_filter]
    rows = rows[:400]  # cap for performance

    entity_counts = {}
    timeline = []
    cooccurrence = {}

    for row in rows:
        text = row["text"]
        found_entities = {e for e in _ENTITY_RE.findall(text) if e not in _STOPWORDS and len(e) > 2}
        for e in found_entities:
            entity_counts[e] = entity_counts.get(e, 0) + 1

        for d in set(_DATE_RE.findall(text)):
            timeline.append({"date": d, "doc_name": row["doc_name"], "page": row["page"],
                              "context": text[:160]})

        ents = list(found_entities)
        for i in range(len(ents)):
            for j in range(i + 1, len(ents)):
                pair = tuple(sorted((ents[i], ents[j])))
                cooccurrence[pair] = cooccurrence.get(pair, 0) + 1

    top_entities = sorted(entity_counts.items(), key=lambda x: -x[1])[:25]
    top_edges = sorted(cooccurrence.items(), key=lambda x: -x[1])[:40]

    graph = {
        "nodes": [{"id": e, "weight": c} for e, c in top_entities],
        "edges": [{"source": p[0], "target": p[1], "weight": w} for p, w in top_edges
                  if p[0] in dict(top_entities) and p[1] in dict(top_entities)],
    }

    timeline.sort(key=lambda x: x["date"])

    return {
        "entities": [{"name": e, "mentions": c} for e, c in top_entities],
        "timeline": timeline[:100],
        "knowledge_graph": graph,
    }

def background_caption_images():
    """Background task to process image captions from the queue."""
    from vision_extractor import extract_via_vision
    global bm25, bm25_documents, bm25_metadata
    while True:
        with caption_queue_lock:
            if not pending_image_captions:
                break
            item = pending_image_captions.pop(0)

        image_path = item["image_path"]
        doc_name = item["doc_name"]
        page = item["page"]
        file_path = item["file_path"]
        
        try:
            res = extract_via_vision(image_path)
            caption = res.get("text", "")
            if not caption:
                continue

            doc_id = hashlib.md5(doc_name.encode()).hexdigest()[:8]
            chunk_id = f"{doc_id}_img_{hashlib.md5(image_path.encode()).hexdigest()[:8]}"
            
            # Embed caption (multimodal if gemini)
            query_embedding = compute_embeddings([caption], task_type="document", titles=[doc_name], image_paths=[image_path])[0]
            
            # Store in ChromaDB
            col = get_or_create_collection()
            
            meta = {
                "doc_name": doc_name,
                "doc_id": doc_id,
                "chunk_index": -1,
                "page": page,
                "total_chunks": -1,
                "previous_chunk": -1,
                "next_chunk": -1,
                "upload_time": time.strftime("%Y-%m-%d %H:%M:%S"),
                "file_hash": "",
                "file_path": file_path,
                "chunk_id": chunk_id,
                "parent_chunk_id": "",
                "source_type": "image",
                "image_path": image_path,
                "image_caption": caption,
                "ocr_engine_used": "",
            }
            col.add(
                embeddings=[query_embedding],
                documents=[caption],
                ids=[chunk_id],
                metadatas=[meta]
            )

            # Update BM25 and image_metadata safely
            with db_lock:
                bm25_documents.append(caption.split())
                bm25_metadata.append({
                    "text": caption,
                    "page": page,
                    "doc_name": doc_name,
                    "chunk_index": -1,
                    "file_path": file_path,
                    "chunk_id": chunk_id,
                    "parent_chunk_id": "",
                    "source_type": "image",
                    "image_path": image_path,
                    "image_caption": caption,
                    "ocr_engine_used": "",
                })
                bm25 = BM25Okapi(bm25_documents)

                # Update existing image_metadata entry's ocr_text so UI can use it
                img_name = os.path.basename(image_path)
                for img in image_metadata:
                    if img["doc_name"] == doc_name and img["page"] == page and img["path"] and img_name in img["path"]:
                        img["ocr_text"] = caption.lower()
                        break

            print(f"[BACKGROUND CAPTION] Indexed image chunk for {image_path}")

        except Exception as e:
            print(f"[BACKGROUND CAPTION] Error processing {image_path}: {e}")

def build_relations(*args, **kwargs):
    """Fallback function for knowledge graph / document relations."""
    return []
